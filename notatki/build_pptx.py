#!/usr/bin/env python3
"""Generate PPTX presentation: Trening Węchowy w warunkach domowych"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
import os

# --- Colors (Offerflow palette) ---
FG = RGBColor(0x17, 0x17, 0x17)       # hsl(0 0% 9%)
MUTED = RGBColor(0x73, 0x73, 0x73)    # hsl(0 0% 45%)
FAINT = RGBColor(0xA3, 0xA3, 0xA3)    # hsl(0 0% 64%)
BG = RGBColor(0xFF, 0xFF, 0xFF)       # white
BG_SEC = RGBColor(0xF7, 0xF7, 0xF7)   # hsl(0 0% 96.5%)
BORDER = RGBColor(0xEB, 0xEB, 0xEB)   # hsl(0 0% 92%)

FONT_NAME = "Inter"

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

# Layout constants
LEFT = Inches(1.2)
TOP_SECTION = Inches(0.6)
TOP_TITLE = Inches(1.0)
TOP_BODY = Inches(1.8)
CONTENT_W = Inches(10.9)


def add_blank_slide():
    layout = prs.slide_layouts[6]  # blank
    return prs.slides.add_slide(layout)


def add_section_label(slide, text):
    txBox = slide.shapes.add_textbox(LEFT, TOP_SECTION, CONTENT_W, Inches(0.35))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text.upper()
    p.font.size = Pt(11)
    p.font.color.rgb = FAINT
    p.font.name = FONT_NAME
    p.font.bold = False
    p.space_after = Pt(0)


def add_title(slide, text, top=None):
    t = top or TOP_TITLE
    txBox = slide.shapes.add_textbox(LEFT, t, CONTENT_W, Inches(0.7))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(32)
    p.font.color.rgb = FG
    p.font.name = FONT_NAME
    p.font.bold = True
    p.space_after = Pt(8)


def add_subtitle(slide, text, top=None):
    t = top or Inches(1.7)
    txBox = slide.shapes.add_textbox(LEFT, t, CONTENT_W, Inches(0.5))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(14)
    p.font.color.rgb = FAINT
    p.font.name = FONT_NAME


def add_body_text(slide, text, top=None, bold_prefix=None):
    t = top or TOP_BODY
    txBox = slide.shapes.add_textbox(LEFT, t, CONTENT_W, Inches(0.6))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    if bold_prefix:
        run_b = p.add_run()
        run_b.text = bold_prefix
        run_b.font.bold = True
        run_b.font.color.rgb = FG
        run_b.font.size = Pt(18)
        run_b.font.name = FONT_NAME
        run_r = p.add_run()
        run_r.text = text
        run_r.font.color.rgb = MUTED
        run_r.font.size = Pt(18)
        run_r.font.name = FONT_NAME
    else:
        p.text = text
        p.font.size = Pt(18)
        p.font.color.rgb = MUTED
        p.font.name = FONT_NAME
    p.line_spacing = Pt(28)
    return txBox


def add_bullets(slide, items, top=None, numbered=False):
    """items: list of (bold_part, rest_text) tuples or plain strings"""
    t = top or TOP_BODY
    txBox = slide.shapes.add_textbox(LEFT, t, CONTENT_W, Inches(4.5))
    tf = txBox.text_frame
    tf.word_wrap = True

    for i, item in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.space_after = Pt(6)
        p.line_spacing = Pt(26)

        if numbered:
            prefix = f"{i+1}. "
            run_n = p.add_run()
            run_n.text = prefix
            run_n.font.size = Pt(17)
            run_n.font.color.rgb = FG
            run_n.font.name = FONT_NAME
            run_n.font.bold = True

        if isinstance(item, tuple):
            bold_part, rest = item
            if not numbered:
                bullet_run = p.add_run()
                bullet_run.text = "  \u2022  "
                bullet_run.font.size = Pt(17)
                bullet_run.font.color.rgb = FAINT
                bullet_run.font.name = FONT_NAME
            run_b = p.add_run()
            run_b.text = bold_part
            run_b.font.bold = True
            run_b.font.color.rgb = FG
            run_b.font.size = Pt(17)
            run_b.font.name = FONT_NAME
            if rest:
                run_r = p.add_run()
                run_r.text = rest
                run_r.font.color.rgb = MUTED
                run_r.font.size = Pt(17)
                run_r.font.name = FONT_NAME
        else:
            if not numbered:
                bullet_run = p.add_run()
                bullet_run.text = "  \u2022  "
                bullet_run.font.size = Pt(17)
                bullet_run.font.color.rgb = FAINT
                bullet_run.font.name = FONT_NAME
            run = p.add_run()
            run.text = item
            run.font.size = Pt(17)
            run.font.color.rgb = MUTED
            run.font.name = FONT_NAME

    return txBox


def add_highlight_box(slide, text, top, bold_prefix=None):
    box = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, LEFT, top, CONTENT_W, Inches(0.9)
    )
    box.fill.solid()
    box.fill.fore_color.rgb = BG_SEC
    box.line.fill.background()
    box.shadow.inherit = False

    tf = box.text_frame
    tf.word_wrap = True
    tf.margin_left = Pt(16)
    tf.margin_right = Pt(16)
    tf.margin_top = Pt(12)
    tf.margin_bottom = Pt(12)
    p = tf.paragraphs[0]
    if bold_prefix:
        run_b = p.add_run()
        run_b.text = bold_prefix
        run_b.font.bold = True
        run_b.font.color.rgb = FG
        run_b.font.size = Pt(16)
        run_b.font.name = FONT_NAME
        run_r = p.add_run()
        run_r.text = text
        run_r.font.color.rgb = MUTED
        run_r.font.size = Pt(16)
        run_r.font.name = FONT_NAME
    else:
        p.text = text
        p.font.size = Pt(16)
        p.font.color.rgb = MUTED
        p.font.name = FONT_NAME
    p.line_spacing = Pt(24)
    return box


def add_h3(slide, text, top):
    txBox = slide.shapes.add_textbox(LEFT, top, CONTENT_W, Inches(0.4))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(20)
    p.font.color.rgb = FG
    p.font.name = FONT_NAME
    p.font.bold = True
    return txBox


def add_table(slide, headers, rows, top, col_widths=None, small=False):
    n_rows = len(rows) + 1
    n_cols = len(headers)
    w = sum(col_widths) if col_widths else CONTENT_W
    tbl_shape = slide.shapes.add_table(n_rows, n_cols, LEFT, top, int(w), Inches(0.4 * n_rows))
    tbl = tbl_shape.table

    if col_widths:
        for i, cw in enumerate(col_widths):
            tbl.columns[i].width = int(cw)

    font_size = Pt(13) if small else Pt(14)
    header_size = Pt(11) if small else Pt(12)

    for i, h in enumerate(headers):
        cell = tbl.cell(0, i)
        cell.text = h
        cell.fill.solid()
        cell.fill.fore_color.rgb = BG_SEC
        for p in cell.text_frame.paragraphs:
            p.font.size = header_size
            p.font.color.rgb = MUTED
            p.font.name = FONT_NAME
            p.font.bold = True

    for r_idx, row in enumerate(rows):
        for c_idx, val in enumerate(row):
            cell = tbl.cell(r_idx + 1, c_idx)
            cell.text = val
            cell.fill.solid()
            cell.fill.fore_color.rgb = BG
            for p in cell.text_frame.paragraphs:
                p.font.size = font_size
                p.font.color.rgb = MUTED
                p.font.name = FONT_NAME

    return tbl_shape


# ===== SLIDES =====

# 1. Tytuł
s = add_blank_slide()
txBox = s.shapes.add_textbox(LEFT, Inches(2.2), CONTENT_W, Inches(1.5))
tf = txBox.text_frame
tf.word_wrap = True
p = tf.paragraphs[0]
p.text = "Trening Węchowy"
p.font.size = Pt(48)
p.font.color.rgb = FG
p.font.name = FONT_NAME
p.font.bold = True
p2 = tf.add_paragraph()
p2.text = "w warunkach domowych"
p2.font.size = Pt(48)
p2.font.color.rgb = FG
p2.font.name = FONT_NAME
p2.font.bold = True
p2.space_before = Pt(0)

txBox2 = s.shapes.add_textbox(LEFT, Inches(4.0), CONTENT_W, Inches(0.5))
tf2 = txBox2.text_frame
p3 = tf2.paragraphs[0]
p3.text = "Opracowanie: Emilia Chodorowska \u00b7 na podstawie kursu Aromapsychologia Anny Bober"
p3.font.size = Pt(14)
p3.font.color.rgb = FAINT
p3.font.name = FONT_NAME


# 2. Wstęp — mechanizm
s = add_blank_slide()
add_section_label(s, "1 \u00b7 Wst\u0119p")
add_title(s, 'Dlaczego Tw\u00f3j nos \u201ezamilk\u0142\u201d?')
add_bullets(s, [
    ("Grypa", " \u2192 obrz\u0119k tkanek fizycznie blokuje dost\u0119p aromat\u00f3w"),
    ("COVID-19", " \u2192 dro\u017cne przewody nosowe, ale wirus atakuje kom\u00f3rki podporowe i gruczo\u0142y Bowmana"),
    ("", "Neurony w\u0119chowe trac\u0105 \u201esystem podtrzymywania \u017cycia\u201d \u2014 jak sprawne odbiorniki bez zasilania"),
    ("", "Brak stymulacji \u2192 atrofia opuszki w\u0119chowej i zmiany w hipokampie"),
])


# 3. Wstęp — nadzieja
s = add_blank_slide()
add_section_label(s, "1 \u00b7 Wst\u0119p")
add_title(s, "Dlaczego to minie?")
add_body_text(s, "Neurony w\u0119chowe maj\u0105 unikaln\u0105 zdolno\u015b\u0107 do regeneracji.")
add_highlight_box(s,
    ' Systematyczny trening w\u0119chowy wykazuje skuteczno\u015b\u0107 por\u00f3wnywaln\u0105 z terapi\u0105 sterydow\u0105. '
    'Regularna stymulacja zwi\u0119ksza obj\u0119to\u015b\u0107 istoty szarej w m\u00f3zgu, odwracaj\u0105c negatywne skutki anosmii. '
    'Tw\u00f3j m\u00f3zg jest plastyczny \u2014 trening to proces jego fizycznej odbudowy.',
    Inches(2.8),
    bold_prefix='\u201e'
)


# 4. Warsztat zapachowy
s = add_blank_slide()
add_section_label(s, "2 \u00b7 Warsztat zapachowy")
add_title(s, "Co przygotowa\u0107?")
add_body_text(s, "Potrzebujemy stworzy\u0107 headspace \u2014 nasycon\u0105 cz\u0105steczkami przestrze\u0144 nad \u017ar\u00f3d\u0142em zapachu.")
add_h3(s, "Niezb\u0119dne wyposa\u017cenie", Inches(2.7))
add_bullets(s, [
    ("S\u0142oiczki z ciemnego szk\u0142a (15-30 ml)", " \u2014 chroni\u0105 olejki, koncentruj\u0105 opary"),
    ("Papier akwarelowy", " \u2014 porowato\u015b\u0107 idealnie trzyma aromat"),
    ("Olejki eteryczne", " \u2014 wy\u0142\u0105cznie naturalne, wysokiej jako\u015bci"),
], top=Inches(3.2))


# 5. Przygotowanie słoiczka
s = add_blank_slide()
add_section_label(s, "2 \u00b7 Warsztat zapachowy")
add_title(s, "Przygotowanie s\u0142oiczka")
add_bullets(s, [
    "W\u0142\u00f3\u017c do s\u0142oiczka pasek papieru akwarelowego",
    "Nas\u0105cz go 4-8 kroplami wybranego olejku",
    "Szczelnie zakr\u0119\u0107, odczekaj godzin\u0119",
    "Co tydzie\u0144 wymieniaj papier i dolewaj olejku (cytrusy szybko oksyduj\u0105)",
    "Popro\u015b kogo\u015b ze sprawnym w\u0119chem o weryfikacj\u0119 intensywno\u015bci",
], numbered=True)


# 6. Wybór zapachów
s = add_blank_slide()
add_section_label(s, "3 \u00b7 Wyb\u00f3r zapach\u00f3w")
add_title(s, "Cztery fundamenty treningu")
add_table(s,
    ["Grupa zapachowa", "Zamienniki", "Dlaczego?"],
    [
        ["Kwiatowa (R\u00f3\u017ca)", "Geranium, ylang-ylang", "Subtelne receptory"],
        ["Owocowa (Cytryna)", "Pomara\u0144cza, grejpfrut", "Wysoka intensywno\u015b\u0107"],
        ["Korzenna (Go\u017adziki)", "Cynamon, wanilia", "Zakotwiczenie w pami\u0119ci"],
        ["\u017bywicza (Eukaliptus)", "Mi\u0119ta, rozmaryn", "Nerw tr\u00f3jdzielny (ch\u0142\u00f3d)"],
    ],
    top=Inches(1.9),
    col_widths=[Inches(4), Inches(3.5), Inches(3.4)]
)
add_highlight_box(s,
    " Wybieraj aromaty budz\u0105ce silne wspomnienia. Emocjonalny \u015blad w m\u00f3zgu u\u0142atwia regeneracj\u0119 po\u0142\u0105cze\u0144 synaptycznych.",
    Inches(4.8),
    bold_prefix="Pami\u0119\u0107 w\u0119chowa:"
)


# 7. Technika małych wdechów
s = add_blank_slide()
add_section_label(s, "4 \u00b7 Technika oddechowa")
add_title(s, 'Technika \u201eMa\u0142ych Wdech\u00f3w\u201d')
add_body_text(s, "G\u0142\u0119boki wdech omija nab\u0142onek w\u0119chowy \u2014 kieruje powietrze prosto do p\u0142uc.")
add_highlight_box(s,
    " Kr\u00f3tkie, ma\u0142e wdechy \u2014 jak pies na spacerze. Tworzysz zawirowania powietrza, "
    "kt\u00f3re kieruj\u0105 headspace bezpo\u015brednio na pole w\u0119chowe.",
    Inches(2.8),
    bold_prefix="Prawid\u0142owa technika:"
)


# 8. Sesja treningowa
s = add_blank_slide()
add_section_label(s, "4 \u00b7 Sesja treningowa")
add_title(s, "Sesja treningowa (~2 min.)")
add_bullets(s, [
    ("☐ ", "Wybierz spokojne miejsce, wycisz telefon"),
    ("☐ ", "Otwórz słoiczek, zbliż go do nosa"),
    ("☐ ", "20 sekund wąchania techniką małych wdechów"),
    ("☐ ", "Zamknij słoiczek \u2014 10-15 sek. przerwy (neutralne powietrze)"),
    ("☐ ", "Przejdź do kolejnego zapachu"),
    ("☐ ", "Powtarzaj 2x dziennie: rano i wieczorem"),
])


# 9. Praca mentalna
s = add_blank_slide()
add_section_label(s, "5 \u00b7 Praca mentalna")
add_title(s, "Wąchanie wyobraźnią")
add_body_text(s, "Nawet przy absolutnej pustce Twoja kora węchowa może wykazywać aktywność.")
add_bullets(s, [
    ("Zamknij oczy", " podczas wąchania"),
    ("Przywołaj obraz obiektu", " \u2014 kolor, teksturę, smak"),
    ("", 'Spróbuj \u201epoczuć\u201d zapach siłą woli'),
    ("Wspieraj się bodźcami wizualnymi", " \u2014 zdjęcia, obrazy"),
], top=Inches(2.8))


# 10. Dzienniczek
s = add_blank_slide()
add_section_label(s, "6 \u00b7 Dzienniczek postępów")
add_title(s, "Cierpliwość i śledzenie postępów")
add_bullets(s, [
    ("Pierwsze efekty:", " zazwyczaj po 4 miesiącach"),
    ("Pełna rehabilitacja:", " 14-24 miesiące"),
], top=Inches(1.9))
add_highlight_box(s,
    " Nieprzyjemne, zniekształcone zapachy (np. spalona guma zamiast kawy) "
    "to dowód, że neurony nawiązują nowe połączenia.",
    Inches(3.2),
    bold_prefix="Parosmia = dobry znak!"
)
add_table(s,
    ["Pole", "Wpis"],
    [
        ["Data", ".................."],
        ["Zapach", ".................."],
        ["Odczucia", "nic / chłód / zniekształcony / czysty"],
        ["Intensywność", "0 \u2013 1 \u2013 2 \u2013 3 \u2013 4 \u2013 5"],
    ],
    top=Inches(4.5),
    col_widths=[Inches(2.5), Inches(5)],
    small=True
)


# 11. Szersze korzyści
s = add_blank_slide()
add_section_label(s, "7 \u00b7 Szersze korzyści")
add_title(s, "Nie tylko po wirusie")
add_body_text(s, "Trening węchowy przynosi szersze korzyści dla mózgu:")
add_bullets(s, [
    ("Poprawa funkcji poznawczych", " \u2014 udowodniona u osób starszych"),
    ("Poprawa płynności semantycznej i werbalnej", ""),
    ("Zwiększenie objętości istoty szarej", " \u2014 odwraca skutki anosmii"),
    ("Wydłużenie życia neuronów węchowych", ""),
    ("Poprawa nastroju", " \u2014 potwierdzona klinicznie"),
], top=Inches(2.8))


# 12. Neuroplastyczność
s = add_blank_slide()
add_section_label(s, "7 \u00b7 Neuroplastyczność")
add_title(s, "Mózg się odbudowuje")
add_h3(s, "Istota szara", Inches(1.8))
add_body_text(s,
    "Anosmia powoduje utratę istoty szarej. Systematyczny trening olfaktoryczny "
    "fizycznie zwiększa jej objętość, odwracając negatywne skutki.",
    top=Inches(2.3)
)
add_h3(s, "Łączność strukturalna", Inches(3.4))
add_body_text(s,
    "Długoterminowa ekspozycja na bodźce węchowe przebudowuje szlaki nerwowe. "
    "Nawet nocna ekspozycja (2h/noc przez 6 miesięcy) poprawia łączność "
    "między układem limbicznym a korą mózgową.",
    top=Inches(3.9)
)


# 13. Złote zasady
s = add_blank_slide()
add_section_label(s, "8 \u00b7 Podsumowanie")
add_title(s, "Złote zasady cierpliwego odkrywcy")
add_bullets(s, [
    ("SYSTEMATYCZNOŚĆ", " \u2014 2x dziennie, codziennie. To Twoje lekarstwo."),
    ("TECHNIKA ODDECHU", ' \u2014 Krótkie, \u201ewęszące\u201d wdechy.'),
    ("WYOBRAŹNIA", " \u2014 Mózg reaguje na wspomnienie zapachu tak samo intensywnie jak na bodziec."),
    ("STYMULACJA TRÓJDZIELNA", " \u2014 Zawsze mięta lub eukaliptus w zestawie."),
    ("CZAS I CIERPLIWOŚĆ", " \u2014 4 miesiące na pierwszy sygnał powrotu."),
], numbered=True)


# ===== SAVE =====
out_path = os.path.join(
    os.path.dirname(os.path.abspath(__file__)),
    "Trening Węchowy — Prezentacja.pptx"
)
prs.save(out_path)
print(f"PPTX: {out_path}")
