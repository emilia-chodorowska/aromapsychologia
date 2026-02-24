#!/usr/bin/env python3
"""Generate PPTX presentation: Trening Węchowy — styl doTERRA (split layout, navy+green)"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
import os

# --- Colors (doTERRA-inspired palette) ---
NAVY = RGBColor(0x1B, 0x3A, 0x4B)
NAVY_LIGHT = RGBColor(0x2C, 0x5F, 0x7C)
TEAL = RGBColor(0x3A, 0x8F, 0x85)
GREEN = RGBColor(0x6B, 0xBF, 0x8A)
GREEN_LIGHT = RGBColor(0xA8, 0xDE, 0xB5)
GREEN_PALE = RGBColor(0xE8, 0xF5, 0xEC)
GRAY_TEXT = RGBColor(0x5A, 0x6A, 0x72)
GRAY_LIGHT = RGBColor(0x8A, 0x96, 0x9C)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)

FONT_NAME = "Inter"

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

# Layout constants — left 55% for content
CONTENT_LEFT = Inches(0.9)
CONTENT_W = Inches(6.3)
PANEL_LEFT = Inches(7.333)
PANEL_W = Inches(6.0)
SLIDE_H = Inches(7.5)


def add_blank_slide():
    layout = prs.slide_layouts[6]  # blank
    return prs.slides.add_slide(layout)


def add_right_panel(slide, color=NAVY):
    panel = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, PANEL_LEFT, Inches(0), PANEL_W, SLIDE_H
    )
    panel.fill.solid()
    panel.fill.fore_color.rgb = color
    panel.line.fill.background()
    panel.shadow.inherit = False
    return panel


def add_section_label(slide, text, top=Inches(0.7)):
    box = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, CONTENT_LEFT, top, Inches(2.8), Inches(0.35)
    )
    box.fill.solid()
    box.fill.fore_color.rgb = GREEN_PALE
    box.line.fill.background()
    box.shadow.inherit = False
    tf = box.text_frame
    tf.margin_left = Pt(10)
    tf.margin_right = Pt(10)
    tf.margin_top = Pt(4)
    tf.margin_bottom = Pt(4)
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.LEFT
    run = p.add_run()
    run.text = text.upper()
    run.font.size = Pt(10)
    run.font.color.rgb = TEAL
    run.font.name = FONT_NAME
    run.font.bold = True


def add_title(slide, text, top=Inches(1.2)):
    txBox = slide.shapes.add_textbox(CONTENT_LEFT, top, CONTENT_W, Inches(0.9))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(36)
    p.font.color.rgb = NAVY
    p.font.name = FONT_NAME
    p.font.bold = True
    p.space_after = Pt(4)
    return txBox


def add_body_text(slide, text, top=Inches(2.1), bold_prefix=None):
    txBox = slide.shapes.add_textbox(CONTENT_LEFT, top, CONTENT_W, Inches(0.7))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    if bold_prefix:
        run_b = p.add_run()
        run_b.text = bold_prefix
        run_b.font.bold = True
        run_b.font.color.rgb = NAVY
        run_b.font.size = Pt(17)
        run_b.font.name = FONT_NAME
        run_r = p.add_run()
        run_r.text = text
        run_r.font.color.rgb = GRAY_TEXT
        run_r.font.size = Pt(17)
        run_r.font.name = FONT_NAME
    else:
        p.text = text
        p.font.size = Pt(17)
        p.font.color.rgb = GRAY_TEXT
        p.font.name = FONT_NAME
    p.line_spacing = Pt(27)
    return txBox


def add_bullets(slide, items, top=Inches(2.1), numbered=False):
    txBox = slide.shapes.add_textbox(CONTENT_LEFT, top, CONTENT_W, Inches(4.5))
    tf = txBox.text_frame
    tf.word_wrap = True

    for i, item in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.space_after = Pt(5)
        p.line_spacing = Pt(25)

        if numbered:
            run_n = p.add_run()
            run_n.text = f"{i+1}. "
            run_n.font.size = Pt(16)
            run_n.font.color.rgb = NAVY
            run_n.font.name = FONT_NAME
            run_n.font.bold = True

        if isinstance(item, tuple):
            bold_part, rest = item
            if not numbered:
                bullet_run = p.add_run()
                bullet_run.text = "  \u2022  "
                bullet_run.font.size = Pt(16)
                bullet_run.font.color.rgb = GREEN
                bullet_run.font.name = FONT_NAME
            run_b = p.add_run()
            run_b.text = bold_part
            run_b.font.bold = True
            run_b.font.color.rgb = NAVY
            run_b.font.size = Pt(16)
            run_b.font.name = FONT_NAME
            if rest:
                run_r = p.add_run()
                run_r.text = rest
                run_r.font.color.rgb = GRAY_TEXT
                run_r.font.size = Pt(16)
                run_r.font.name = FONT_NAME
        else:
            if not numbered:
                bullet_run = p.add_run()
                bullet_run.text = "  \u2022  "
                bullet_run.font.size = Pt(16)
                bullet_run.font.color.rgb = GREEN
                bullet_run.font.name = FONT_NAME
            run = p.add_run()
            run.text = item
            run.font.size = Pt(16)
            run.font.color.rgb = GRAY_TEXT
            run.font.name = FONT_NAME

    return txBox


def add_green_box(slide, text, top, bold_prefix=None):
    box = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, CONTENT_LEFT, top, CONTENT_W, Inches(0.85)
    )
    box.fill.solid()
    box.fill.fore_color.rgb = GREEN
    box.line.fill.background()
    box.shadow.inherit = False

    tf = box.text_frame
    tf.word_wrap = True
    tf.margin_left = Pt(14)
    tf.margin_right = Pt(14)
    tf.margin_top = Pt(10)
    tf.margin_bottom = Pt(10)
    p = tf.paragraphs[0]
    if bold_prefix:
        run_b = p.add_run()
        run_b.text = bold_prefix
        run_b.font.bold = True
        run_b.font.color.rgb = WHITE
        run_b.font.size = Pt(15)
        run_b.font.name = FONT_NAME
        run_r = p.add_run()
        run_r.text = text
        run_r.font.color.rgb = WHITE
        run_r.font.size = Pt(15)
        run_r.font.name = FONT_NAME
    else:
        p.text = text
        p.font.size = Pt(15)
        p.font.color.rgb = WHITE
        p.font.name = FONT_NAME
    p.line_spacing = Pt(23)
    return box


def add_navy_box(slide, text, top, bold_prefix=None):
    box = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, CONTENT_LEFT, top, CONTENT_W, Inches(0.85)
    )
    box.fill.solid()
    box.fill.fore_color.rgb = NAVY
    box.line.fill.background()
    box.shadow.inherit = False

    tf = box.text_frame
    tf.word_wrap = True
    tf.margin_left = Pt(14)
    tf.margin_right = Pt(14)
    tf.margin_top = Pt(10)
    tf.margin_bottom = Pt(10)
    p = tf.paragraphs[0]
    if bold_prefix:
        run_b = p.add_run()
        run_b.text = bold_prefix
        run_b.font.bold = True
        run_b.font.color.rgb = WHITE
        run_b.font.size = Pt(15)
        run_b.font.name = FONT_NAME
        run_r = p.add_run()
        run_r.text = text
        run_r.font.color.rgb = RGBColor(0xCC, 0xDD, 0xDD)
        run_r.font.size = Pt(15)
        run_r.font.name = FONT_NAME
    else:
        p.text = text
        p.font.size = Pt(15)
        p.font.color.rgb = RGBColor(0xCC, 0xDD, 0xDD)
        p.font.name = FONT_NAME
    p.line_spacing = Pt(23)
    return box


def add_h3(slide, text, top):
    txBox = slide.shapes.add_textbox(CONTENT_LEFT, top, CONTENT_W, Inches(0.4))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(19)
    p.font.color.rgb = NAVY
    p.font.name = FONT_NAME
    p.font.bold = True
    return txBox


def add_table(slide, headers, rows, top, col_widths=None, small=False):
    n_rows = len(rows) + 1
    n_cols = len(headers)
    w = sum(col_widths) if col_widths else int(CONTENT_W)
    tbl_shape = slide.shapes.add_table(n_rows, n_cols, CONTENT_LEFT, top, int(w), Inches(0.38 * n_rows))
    tbl = tbl_shape.table

    if col_widths:
        for i, cw in enumerate(col_widths):
            tbl.columns[i].width = int(cw)

    font_size = Pt(12) if small else Pt(13)
    header_size = Pt(11) if small else Pt(12)

    for i, h in enumerate(headers):
        cell = tbl.cell(0, i)
        cell.text = h
        cell.fill.solid()
        cell.fill.fore_color.rgb = NAVY
        for p in cell.text_frame.paragraphs:
            p.font.size = header_size
            p.font.color.rgb = WHITE
            p.font.name = FONT_NAME
            p.font.bold = True

    for r_idx, row in enumerate(rows):
        for c_idx, val in enumerate(row):
            cell = tbl.cell(r_idx + 1, c_idx)
            cell.text = val
            cell.fill.solid()
            cell.fill.fore_color.rgb = GREEN_PALE if r_idx % 2 == 0 else WHITE
            for p in cell.text_frame.paragraphs:
                p.font.size = font_size
                p.font.color.rgb = GRAY_TEXT
                p.font.name = FONT_NAME

    return tbl_shape


# ===== SLIDES =====

# 1. Tytuł
s = add_blank_slide()
add_right_panel(s, NAVY)
txBox = s.shapes.add_textbox(CONTENT_LEFT, Inches(2.0), CONTENT_W, Inches(1.8))
tf = txBox.text_frame
tf.word_wrap = True
p = tf.paragraphs[0]
p.text = "Trening W\u0119chowy"
p.font.size = Pt(48)
p.font.color.rgb = NAVY
p.font.name = FONT_NAME
p.font.bold = True
p2 = tf.add_paragraph()
p2.text = "w warunkach domowych"
p2.font.size = Pt(48)
p2.font.color.rgb = NAVY
p2.font.name = FONT_NAME
p2.font.bold = True
p2.space_before = Pt(0)

txBox2 = s.shapes.add_textbox(CONTENT_LEFT, Inches(4.2), CONTENT_W, Inches(0.5))
tf2 = txBox2.text_frame
p3 = tf2.paragraphs[0]
p3.text = "Opracowanie: Emilia Chodorowska \u00b7 na podstawie kursu Aromapsychologia Anny Bober"
p3.font.size = Pt(13)
p3.font.color.rgb = GRAY_LIGHT
p3.font.name = FONT_NAME


# 2. Wstęp — mechanizm
s = add_blank_slide()
add_right_panel(s, NAVY)
add_section_label(s, "1 \u00b7 Wst\u0119p")
add_title(s, 'Dlaczego Tw\u00f3j nos \u201ezamilk\u0142\u201d?')
add_bullets(s, [
    ("Grypa", " \u2192 obrz\u0119k tkanek blokuje dost\u0119p aromat\u00f3w"),
    ("COVID-19", " \u2192 dro\u017cne przewody nosowe, ale wirus atakuje kom\u00f3rki podporowe i gruczo\u0142y Bowmana"),
    ("", "Neurony trac\u0105 \u201esystem podtrzymywania \u017cycia\u201d \u2014 jak odbiorniki bez zasilania"),
    ("", "Brak stymulacji \u2192 atrofia opuszki w\u0119chowej i zmiany w hipokampie"),
])


# 3. Wstęp — nadzieja
s = add_blank_slide()
add_right_panel(s, TEAL)
add_section_label(s, "1 \u00b7 Wst\u0119p")
add_title(s, "Dlaczego to minie?")
add_body_text(s, "Neurony w\u0119chowe maj\u0105 unikaln\u0105 zdolno\u015b\u0107 do regeneracji.")
add_green_box(s,
    ' Systematyczny trening w\u0119chowy wykazuje skuteczno\u015b\u0107 por\u00f3wnywaln\u0105 z terapi\u0105 sterydow\u0105. '
    'Regularna stymulacja zwi\u0119ksza obj\u0119to\u015b\u0107 istoty szarej w m\u00f3zgu. '
    'Tw\u00f3j m\u00f3zg jest plastyczny \u2014 trening to proces jego fizycznej odbudowy.',
    Inches(3.0),
    bold_prefix='\u201e'
)


# 4. Warsztat zapachowy
s = add_blank_slide()
add_right_panel(s, GREEN)
add_section_label(s, "2 \u00b7 Warsztat zapachowy")
add_title(s, "Co przygotowa\u0107?")
add_body_text(s, "Potrzebujemy stworzy\u0107 headspace \u2014 nasycon\u0105 cz\u0105steczkami przestrze\u0144 nad \u017ar\u00f3d\u0142em zapachu.")
add_h3(s, "Niezb\u0119dne wyposa\u017cenie", Inches(2.9))
add_bullets(s, [
    ("S\u0142oiczki z ciemnego szk\u0142a (15-30 ml)", " \u2014 chroni\u0105 olejki, koncentruj\u0105 opary"),
    ("Papier akwarelowy", " \u2014 porowato\u015b\u0107 idealnie trzyma aromat"),
    ("Olejki eteryczne", " \u2014 wy\u0142\u0105cznie naturalne, wysokiej jako\u015bci"),
], top=Inches(3.4))


# 5. Przygotowanie słoiczka
s = add_blank_slide()
add_right_panel(s, NAVY_LIGHT)
add_section_label(s, "2 \u00b7 Warsztat zapachowy")
add_title(s, "Przygotowanie s\u0142oiczka")
add_bullets(s, [
    "W\u0142\u00f3\u017c do s\u0142oiczka pasek papieru akwarelowego",
    "Nas\u0105cz go 4-8 kroplami wybranego olejku",
    "Szczelnie zakr\u0119\u0107, odczekaj godzin\u0119",
    "Co tydzie\u0144 wymieniaj papier i dolewaj olejku (cytrusy szybko oksyduj\u0105)",
    "Popro\u015b kogo\u015b ze sprawnym w\u0119chem o weryfikacj\u0119 intensywno\u015bci",
], numbered=True)


# 6. Wybór zapachów
s = add_blank_slide()
add_right_panel(s, TEAL)
add_section_label(s, "3 \u00b7 Wyb\u00f3r zapach\u00f3w")
add_title(s, "Cztery fundamenty treningu")
add_table(s,
    ["Grupa zapachowa", "Zamienniki", "Dlaczego?"],
    [
        ["Kwiatowa (R\u00f3\u017ca)", "Geranium, ylang-ylang", "Subtelne receptory"],
        ["Owocowa (Cytryna)", "Pomara\u0144cza, grejpfrut", "Wysoka intensywno\u015b\u0107"],
        ["Korzenna (Go\u017adziki)", "Cynamon, wanilia", "Zakotwiczenie w pami\u0119ci"],
        ["\u017bywicza (Eukaliptus)", "Mi\u0119ta, rozmaryn", "Nerw tr\u00f3jdzielny (ch\u0142\u00f3d)"],
    ],
    top=Inches(2.0),
    col_widths=[Inches(2.5), Inches(2.1), Inches(2.1)]
)
add_green_box(s,
    " Wybieraj aromaty budz\u0105ce silne wspomnienia. Emocjonalny \u015blad u\u0142atwia regeneracj\u0119.",
    Inches(4.8),
    bold_prefix="Pami\u0119\u0107 w\u0119chowa:"
)


# 7. Technika małych wdechów
s = add_blank_slide()
add_right_panel(s, GREEN)
add_section_label(s, "4 \u00b7 Technika oddechowa")
add_title(s, 'Technika \u201eMa\u0142ych Wdech\u00f3w\u201d')
add_body_text(s, "G\u0142\u0119boki wdech omija nab\u0142onek w\u0119chowy \u2014 kieruje powietrze prosto do p\u0142uc.")
add_navy_box(s,
    " Kr\u00f3tkie, ma\u0142e wdechy \u2014 jak pies na spacerze. Tworzysz zawirowania powietrza, "
    "kt\u00f3re kieruj\u0105 headspace bezpo\u015brednio na pole w\u0119chowe.",
    Inches(3.0),
    bold_prefix="Prawid\u0142owa technika:"
)


# 8. Sesja treningowa
s = add_blank_slide()
add_right_panel(s, NAVY)
add_section_label(s, "4 \u00b7 Sesja treningowa")
add_title(s, "Sesja treningowa (~2 min.)")
add_bullets(s, [
    ("\u2610 ", "Wybierz spokojne miejsce, wycisz telefon"),
    ("\u2610 ", "Otw\u00f3rz s\u0142oiczek, zbli\u017c go do nosa"),
    ("\u2610 ", "20 sekund w\u0105chania technik\u0105 ma\u0142ych wdech\u00f3w"),
    ("\u2610 ", "Zamknij s\u0142oiczek \u2014 10-15 sek. przerwy"),
    ("\u2610 ", "Przejd\u017a do kolejnego zapachu"),
    ("\u2610 ", "Powtarzaj 2x dziennie: rano i wieczorem"),
])


# 9. Praca mentalna
s = add_blank_slide()
add_right_panel(s, GREEN)
add_section_label(s, "5 \u00b7 Praca mentalna")
add_title(s, "W\u0105chanie wyobra\u017ani\u0105")
add_body_text(s, "Nawet przy absolutnej pustce kora w\u0119chowa mo\u017ce wykazywa\u0107 aktywno\u015b\u0107.")
add_bullets(s, [
    ("Zamknij oczy", " podczas w\u0105chania"),
    ("Przywo\u0142aj obraz obiektu", " \u2014 kolor, tekstur\u0119, smak"),
    ("", 'Spr\u00f3buj \u201epoczu\u0107\u201d zapach si\u0142\u0105 woli'),
    ("Wspieraj si\u0119 bod\u017acami wizualnymi", " \u2014 zdj\u0119cia, obrazy"),
], top=Inches(3.0))
add_green_box(s,
    "Medytacja sensoryczna zapobiega degradacji neuron\u00f3w i stymuluje je do d\u0142u\u017cszego prze\u017cycia.",
    Inches(5.2)
)


# 10. Dzienniczek
s = add_blank_slide()
add_right_panel(s, NAVY_LIGHT)
add_section_label(s, "6 \u00b7 Dzienniczek post\u0119p\u00f3w")
add_title(s, "Cierpliwo\u015b\u0107 i \u015bledzenie post\u0119p\u00f3w")
add_bullets(s, [
    ("Pierwsze efekty:", " zazwyczaj po 4 miesi\u0105cach"),
    ("Pe\u0142na rehabilitacja:", " 14-24 miesi\u0119cy"),
], top=Inches(2.1))
add_green_box(s,
    " Nieprzyjemne, zniekszta\u0142cone zapachy to dow\u00f3d, \u017ce neurony nawi\u0105zuj\u0105 nowe po\u0142\u0105czenia.",
    Inches(3.3),
    bold_prefix="Parosmia = dobry znak!"
)
add_table(s,
    ["Pole", "Wpis"],
    [
        ["Data", ".................."],
        ["Zapach", ".................."],
        ["Odczucia", "nic / ch\u0142\u00f3d / zniekszta\u0142cony / czysty"],
        ["Intensywno\u015b\u0107", "0 \u2013 1 \u2013 2 \u2013 3 \u2013 4 \u2013 5"],
    ],
    top=Inches(4.5),
    col_widths=[Inches(2.2), Inches(4.3)],
    small=True
)


# 11. Szersze korzyści
s = add_blank_slide()
add_right_panel(s, TEAL)
add_section_label(s, "7 \u00b7 Szersze korzy\u015bci")
add_title(s, "Nie tylko po wirusie")
add_body_text(s, "Trening w\u0119chowy przynosi szersze korzy\u015bci dla m\u00f3zgu:")
add_bullets(s, [
    ("Poprawa funkcji poznawczych", " \u2014 udowodniona u os\u00f3b starszych"),
    ("Poprawa p\u0142ynno\u015bci semantycznej i werbalnej", ""),
    ("Zwi\u0119kszenie obj\u0119to\u015bci istoty szarej", " \u2014 odwraca skutki anosmii"),
    ("Wyd\u0142u\u017cenie \u017cycia neuron\u00f3w w\u0119chowych", ""),
    ("Poprawa nastroju", " \u2014 potwierdzona klinicznie"),
], top=Inches(3.0))


# 12. Neuroplastyczność
s = add_blank_slide()
add_right_panel(s, NAVY)
add_section_label(s, "7 \u00b7 Neuroplastyczno\u015b\u0107")
add_title(s, "M\u00f3zg si\u0119 odbudowuje")
add_h3(s, "Istota szara", Inches(2.0))
add_body_text(s,
    "Anosmia powoduje utrat\u0119 istoty szarej. Systematyczny trening "
    "fizycznie zwi\u0119ksza jej obj\u0119to\u015b\u0107, odwracaj\u0105c negatywne skutki.",
    top=Inches(2.5)
)
add_h3(s, "\u0141\u0105czno\u015b\u0107 strukturalna", Inches(3.6))
add_body_text(s,
    "D\u0142ugoterminowa ekspozycja na bod\u017ace w\u0119chowe przebudowuje szlaki nerwowe. "
    "Nawet nocna ekspozycja (2h/noc przez 6 miesi\u0119cy) poprawia \u0142\u0105czno\u015b\u0107 "
    "mi\u0119dzy uk\u0142adem limbicznym a kor\u0105 m\u00f3zgow\u0105.",
    top=Inches(4.1)
)
add_green_box(s,
    "Trening w\u0119chowy to tak\u017ce trening umys\u0142u i pami\u0119ci dla senior\u00f3w.",
    Inches(5.3)
)


# 13. Złote zasady
s = add_blank_slide()
add_right_panel(s, NAVY)
add_section_label(s, "8 \u00b7 Podsumowanie")
add_title(s, "Z\u0142ote zasady cierpliwego odkrywcy")
add_bullets(s, [
    ("SYSTEMATYCZNO\u015a\u0106", " \u2014 2x dziennie, codziennie. To Twoje lekarstwo."),
    ("TECHNIKA ODDECHU", ' \u2014 Kr\u00f3tkie, \u201ew\u0119sz\u0105ce\u201d wdechy.'),
    ("WYOBRA\u0179NIA", " \u2014 M\u00f3zg reaguje na wspomnienie zapachu tak samo intensywnie."),
    ("STYMULACJA TR\u00d3JDZIELNA", " \u2014 Zawsze mi\u0119ta lub eukaliptus w zestawie."),
    ("CZAS I CIERPLIWO\u015a\u0106", " \u2014 4 miesi\u0105ce na pierwszy sygna\u0142 powrotu."),
], numbered=True)


# ===== SAVE =====
out_path = os.path.join(
    os.path.dirname(os.path.abspath(__file__)),
    "Trening W\u0119chowy \u2014 Prezentacja \u2014 doTERRA.pptx"
)
prs.save(out_path)
print(f"PPTX: {out_path}")
